import os
import json
import io
import time
from datetime import datetime
from django.utils import timezone
from django.core.management.base import BaseCommand
from openpyxl import Workbook
from pptx import Presentation
import fitz  # PyMuPDF
from microsoft.models import Email, EmailAccount
from ai_orchestrator.services.ai_processor import AIProcessorService
from ai_orchestrator.services.document_processor import DocumentProcessorService

class Command(BaseCommand):
    help = 'Stress test the AI pipeline with 5 complex, multi-attachment emails'

    def handle(self, *args, **options):
        account, _ = EmailAccount.objects.get_or_create(email="stress-test@india-alt.com")
        doc_processor = DocumentProcessorService()
        ai_service = AIProcessorService()
        
        scenarios = [
            {
                "id": 1,
                "name": "Scenario 1: Large Text-Only Deal",
                "subject": "Confidential: Project Maharaja - Consumer Goods",
                "body": "This is a massive summary. " + ("The company is growing at 40% CAGR. " * 500),
                "attachments": []
            },
            {
                "id": 2,
                "name": "Scenario 2: Complex Excel (Financials)",
                "subject": "Financial Model - HealthTech Innovators FY25",
                "body": "PFA the detailed unit economics for our 200 clinics.",
                "attachments": [("Financial_Model.xlsx", "excel")]
            },
            {
                "id": 3,
                "name": "Scenario 3: Multi-Page PDF (IM)",
                "subject": "Information Memorandum: AquaPure Systems",
                "body": "Attached is the IM for AquaPure.",
                "attachments": [("Information_Memorandum.pdf", "pdf")]
            },
            {
                "id": 4,
                "name": "Scenario 4: Vision & PPT (Strategy)",
                "subject": "Expansion Strategy - RetailForce India",
                "body": "Slide deck attached. Note the charts on page 4.",
                "attachments": [("Strategy_Deck.pptx", "pptx")]
            },
            {
                "id": 5,
                "name": "Scenario 5: TOTAL STRESS (All Types)",
                "subject": "URGENT: Project Omega - Full Diligence Pack",
                "body": "Everything is here for tomorrow's IC meeting.",
                "attachments": [
                    ("Diligence_Pack.pdf", "pdf"),
                    ("CapTable.xlsx", "excel"),
                    ("FinalPitch.pptx", "pptx")
                ]
            }
        ]

        def generate_attachment(file_type, name):
            if file_type == "excel":
                wb = Workbook()
                ws = wb.active
                ws.append(["Month", "Revenue", "COGS", "Marketing", "EBITDA", "Net Profit"])
                for i in range(1, 101):
                    ws.append([f"Month {i}", i*100000, i*40000, i*15000, i*30000, i*25000])
                buf = io.BytesIO()
                wb.save(buf)
                return buf.getvalue()
            elif file_type == "pdf":
                doc = fitz.open()
                for p in range(5):
                    page = doc.new_page()
                    # Fixed f-string syntax
                    msg = f"PAGE {p+1}: Market Analysis for {name}. " + ("Detailed market trends. " * 100)
                    page.insert_text((50, 50), msg)
                return doc.tobytes()
            elif file_type == "pptx":
                prs = Presentation()
                for s in range(5):
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = f"Strategy Slide {s+1}"
                    slide.placeholders[1].text = "Competitive Moat: Network effects. " * 10
                buf = io.BytesIO()
                prs.save(buf)
                return buf.getvalue()
            return b"Mock data"

        self.stdout.write(self.style.MIGRATE_HEADING("\n=== STARTING STRESS TEST (T4 GPU) ===\n"))

        for s in scenarios:
            self.stdout.write(self.style.MIGRATE_LABEL(f"Running {s['name']}..."))
            
            full_context = s['body']
            for filename, ftype in s['attachments']:
                self.stdout.write(f"  - Ingesting {filename}...")
                content = generate_attachment(ftype, s['name'])
                text = doc_processor.extract_text(content, filename)
                full_context += f"\n\nFILE: {filename}\n{text}"

            self.stdout.write(f"  - Data Size: {len(full_context)} chars. Analyzing...")
            
            start_time = time.time()
            result = ai_service.process_content(
                content=full_context,
                skill_name="deal_extraction",
                metadata={'subject': s['subject']},
                source_type="stress_test"
            )
            duration = time.time() - start_time

            if "error" in result:
                self.stdout.write(self.style.ERROR(f"  FAILED in {duration:.2f}s: {result.get('error')}"))
            else:
                self.stdout.write(self.style.SUCCESS(f"  SUCCESS in {duration:.2f}s!"))
                self.stdout.write("\n--- RAW LLM RESPONSE ---")
                self.stdout.write(result.get("_raw_response", "N/A"))
                self.stdout.write("\n--- PARSED JSON ---")
                self.stdout.write(json.dumps(result, indent=2))
                self.stdout.write("\n--- END OUTPUT ---\n")

        self.stdout.write(self.style.MIGRATE_HEADING("\n=== STRESS TEST COMPLETE ===\n"))
