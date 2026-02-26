import os
import json
import io
from datetime import datetime
from django.utils import timezone
from django.core.management.base import BaseCommand
from openpyxl import Workbook
from pptx import Presentation
import fitz  # PyMuPDF
from microsoft.models import Email, EmailAccount
from ai_orchestrator.services.ai_processor import AIProcessorService
from ai_orchestrator.services.document_processor import DocumentProcessorService

class Command(BaseCommand):
    help = 'Tests the AI analysis with PDF, Excel, and PPTX attachments'

    def handle(self, *args, **options):
        # 1. Setup mock account
        account, _ = EmailAccount.objects.get_or_create(email="test-attachment-all@india-alt.com")
        doc_processor = DocumentProcessorService()
        ai_service = AIProcessorService()

        # 2. Generate Files
        wb = Workbook()
        ws = wb.active
        ws.append(["Metric", "Value"])
        ws.append(["Revenue FY24", "250 Cr"])
        excel_bytes = io.BytesIO()
        wb.save(excel_bytes)
        
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[0])
        slide.shapes.title.text = "CloudNine SaaS Pitch"
        pptx_bytes = io.BytesIO()
        prs.save(pptx_bytes)

        pdf_doc = fitz.open()
        page = pdf_doc.new_page()
        page.insert_text((50, 50), "Project Blue: Seeking 25 Cr for fleet electrification.")
        pdf_bytes = pdf_doc.tobytes()
        pdf_doc.close()

        files = [
            ("Financials.xlsx", excel_bytes.getvalue()),
            ("PitchDeck.pptx", pptx_bytes.getvalue()),
            ("OnePager.pdf", pdf_bytes)
        ]

        full_context = "Review the following deal docs:\n"

        for filename, content in files:
            self.stdout.write(f"Processing {filename}...")
            try:
                text = doc_processor.extract_text(content, filename)
                # Hard truncation for test
                full_context += f"\nFILE {filename}: {text[:500]}\n"
            except Exception as e:
                self.stdout.write(self.style.ERROR(f"  Failed {filename}: {str(e)}"))

        # 5. AI Analysis
        self.stdout.write("\nSending to gemma3:4b...")
        result = ai_service.process_content(
            content=full_context,
            skill_name="deal_extraction",
            metadata={'subject': 'Project CloudNine'},
            source_type="multi_file_test"
        )

        self.stdout.write(self.style.SUCCESS("\n--- MULTI-FILE AI RESULT ---"))
        self.stdout.write(json.dumps(result, indent=2))
